from flask import Flask, render_template, url_for,request,redirect
import ollama
from pptx import Presentation
import os
from PIL import Image
from pdf2image import convert_from_path
import subprocess
from spire.presentation.common import *  
from spire.presentation import * 
from createppt import Ppt

app = Flask(__name__)

def generating_ppt(slide, prompt):
    p1 = Presentation()
    #create title
    desired_model = 'llama2:latest'
    def title_for_content(prompt):
        response = ollama.chat(model=desired_model, messages=[
            {
                'role' : 'user',
                'content' : prompt,
            },    
        ])
        return(response['message']['content'])

    #create content
    def content_on_title(title):
        response_content = ollama.chat(model=desired_model, messages=[
            {
                'role' : 'user',
                'content' : f'''Write content on the topic "{title}" in exactly with 100 words. Do not include phrases like "Sure, here is the content" or mention the word count or any introductory lines. Start directly with the content.. '''
            }
        ])
        return(response_content['message']['content'])

    #alter title after creation
    def alter_title(title, dele):
        car_list = [line for line in title if line]
        del car_list[dele]
        return f'''{car_list}'''

    #creating ppt
    def create_ppt(titles_raw):
        for title in titles_raw:
            content = content_on_title(title)
            slide_register = p1.slide_layouts[1]
            slide = p1.slides.add_slide(slide_register)
            slide_title = slide.shapes.title
            slide_content = slide.placeholders[1]
            slide_title.text = f'''{title}'''
            slide_content.text = f'''{content}'''

        direction=r"D:\Code\Tot_PPT\static\ppt_temp"
        path_dir = os.path.join(direction, 'newppt.pptx')   
        p1.save(path_dir)
        

    #altering ppt after creation
    def alter_ppt():
        read_ppt = Presentation("zpresentation.pptx")
        
    #main function
    def main(title,slides):
        headings = title_for_content(prompt=f'''"give me {slides} titles for the keyword from prompt '{title}'. Without any additional text or explanations . Respond with just the titles.''')
        #normal generating heading
        nor_heading = f'''{headings}'''.splitlines()
        filter_heading = [title for title in nor_heading if title]
        # #checking if user need change
        # user_title_stis = int(input("title to remove: ")) 
        # if user_title_stis:
        #     altered_title= alter_title(filter_heading, dele=user_title_stis)
        #     create_ppt(titles_raw=altered_title)

        # #default ppt creation
        # else:
        create_ppt(titles_raw = filter_heading)
        return f"created ppt"

    # main(title=prompt, slides = slide)



#main func
@app.route('/')
def home():
    return render_template('index.html')

#login page
@app.route('/login', methods = ["GET", "POST"])
def login_page():
    if request.method == "POST":
        no_of_slide = request.form.get('email')
        print(no_of_slide)
        
    return render_template('login.html')

#option page
@app.route('/options')
def select_option():
    return render_template('options.html')

#chosing no of slides etc
@app.route('/choice')
def choice_page():

    return render_template('choice.html')

#shows slides
@app.route('/output', methods = ["GET", "POST"])
def ppt():
    if request.method == "POST":
        p = Ppt
        prompt= request.form.get('userprompt')
        slide = request.form.get('noslides')

        generating_ppt(slide, prompt)

    return render_template('output.html')

@app.route('/heading', methods= ['GET', 'POST'])
def headings():
    if request.method == 'POST':
        name = request.form.get('name')
        print(name)
    return render_template('heading.html')

@app.route('/result', methods = ['GET', 'POST'])
def show_presentation():
    ppt_path = 'D:/Code/Tot_PPT/static/ppt_temp/new.pptx'
    output_folder = 'D:/Code/Tot_PPT/static/slides'
    presentation = Presentation()  
    if request.method == "POST":
        prompt= request.form.get('userprompt')
        no_slide = request.form.get('noslides')
        create_ppt = Ppt
        create_ppt.getppt(prompt,no_slide)   ###turn on first
    else:
        print('no input were taken')
  
    presentation.LoadFromFile(ppt_path)
    s = 0
    output_folder = "D:/Code/Tot_PPT/static/slides"
    
    # Ensure the directory exists
    if not os.path.exists(output_folder):
        os.makedirs(output_folder)
    
    # Loop through the slides and save each as an image
    for i, slide in enumerate(presentation.Slides):  
        s += 1
        fileName = f"ToImage{s}.jpg"
        path_dir = os.path.join(output_folder, fileName) 
        
        # Save each slide as an image
        image = slide.SaveAsImage() 
        image.Save(path_dir)  
        image.Dispose()  
    
    presentation.Dispose()

    # Get the list of slide images
    slides = os.listdir(output_folder)
    slides = [f'/static/slides/{slide}' for slide in sorted(slides)]

    # Pass the slide images to the frontend template
    return render_template('testfinal.html', slides=slides)


if __name__ == '__main__':
    app.run(debug=True)
